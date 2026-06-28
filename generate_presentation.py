"""
Acids, Bases & Salts - Class 10 NCERT Complete Presentation Generator
Generates a comprehensive, creative PowerPoint presentation using python-pptx

Installation: pip install python-pptx pillow

Usage: python generate_presentation.py
Output: Acids_Bases_Salts_Class10.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Color Definitions
ACID_COLOR = RGBColor(220, 50, 50)           # Red for acids
BASE_COLOR = RGBColor(50, 100, 220)          # Blue for bases
SALT_COLOR = RGBColor(180, 50, 180)          # Purple for salts
NEUTRAL_COLOR = RGBColor(100, 100, 100)      # Gray for neutral
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
LIGHT_GRAY = RGBColor(240, 240, 240)
DARK_GRAY = RGBColor(50, 50, 50)

def create_presentation():
    """Create the main presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    return prs

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 80, 150)  # Dark blue gradient effect
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "ACIDS, BASES & SALTS"
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "A Journey Through Chemistry"
    subtitle_p.font.size = Pt(44)
    subtitle_p.font.italic = True
    subtitle_p.font.color.rgb = RGBColor(200, 220, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Class info
    class_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    class_frame = class_box.text_frame
    class_p = class_frame.paragraphs[0]
    class_p.text = "Class X | Chemistry | NCERT Complete Course"
    class_p.font.size = Pt(28)
    class_p.font.color.rgb = RGBColor(180, 200, 255)
    class_p.alignment = PP_ALIGN.CENTER

def add_objectives_slide(prs):
    """Slide 2: Learning Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Learning Objectives"
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = DARK_GRAY
    
    objectives = [
        "✓ Understand the properties of acids and bases",
        "✓ Distinguish between acids, bases, and salts",
        "✓ Learn about pH scale and pH indicators",
        "✓ Understand neutralization reactions",
        "✓ Recognize common acids, bases, and salts",
        "✓ Learn about strong and weak acids/bases"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, obj in enumerate(objectives):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = obj
        p.font.size = Pt(26)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_acids_slide(prs):
    """Slide 3: What are Acids?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 230, 230)  # Light red
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "What are Acids?"
    title_p.font.size = Pt(50)
    title_p.font.bold = True
    title_p.font.color.rgb = ACID_COLOR
    
    # Definition
    def_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.8))
    def_frame = def_box.text_frame
    def_frame.word_wrap = True
    def_p = def_frame.paragraphs[0]
    def_p.text = "Definition: Substances that produce H⁺ ions (or H₃O⁺) in aqueous solution"
    def_p.font.size = Pt(20)
    def_p.font.bold = True
    def_p.font.color.rgb = ACID_COLOR
    
    # Characteristics
    char_title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(0.4))
    char_frame = char_title.text_frame
    char_p = char_frame.paragraphs[0]
    char_p.text = "Characteristics:"
    char_p.font.size = Pt(22)
    char_p.font.bold = True
    char_p.font.color.rgb = ACID_COLOR
    
    characteristics = [
        "• Sour taste",
        "• Turn blue litmus paper RED",
        "• Conduct electricity",
        "• React with bases",
        "• React with metals → H₂ gas"
    ]
    
    char_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(4), Inches(3.5))
    char_content_frame = char_content.text_frame
    char_content_frame.word_wrap = True
    
    for i, char in enumerate(characteristics):
        if i > 0:
            char_content_frame.add_paragraph()
        p = char_content_frame.paragraphs[i]
        p.text = char
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
    
    # Examples
    ex_title = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.3), Inches(0.4))
    ex_frame = ex_title.text_frame
    ex_p = ex_frame.paragraphs[0]
    ex_p.text = "Common Examples:"
    ex_p.font.size = Pt(22)
    ex_p.font.bold = True
    ex_p.font.color.rgb = ACID_COLOR
    
    examples = [
        "• HCl - Stomach acid",
        "• H₂SO₄ - Battery acid",
        "• HNO₃ - Fertilizers",
        "• CH₃COOH - Vinegar",
        "• C₆H₈O₇ - Citrus fruits"
    ]
    
    ex_content = slide.shapes.add_textbox(Inches(5.2), Inches(2.5), Inches(4.3), Inches(3.5))
    ex_content_frame = ex_content.text_frame
    ex_content_frame.word_wrap = True
    
    for i, ex in enumerate(examples):
        if i > 0:
            ex_content_frame.add_paragraph()
        p = ex_content_frame.paragraphs[i]
        p.text = ex
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

def add_bases_slide(prs):
    """Slide 4: What are Bases?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 240, 255)  # Light blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "What are Bases?"
    title_p.font.size = Pt(50)
    title_p.font.bold = True
    title_p.font.color.rgb = BASE_COLOR
    
    # Definition
    def_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.8))
    def_frame = def_box.text_frame
    def_frame.word_wrap = True
    def_p = def_frame.paragraphs[0]
    def_p.text = "Definition: Substances that produce OH⁻ ions in aqueous solution"
    def_p.font.size = Pt(20)
    def_p.font.bold = True
    def_p.font.color.rgb = BASE_COLOR
    
    # Characteristics
    char_title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(0.4))
    char_frame = char_title.text_frame
    char_p = char_frame.paragraphs[0]
    char_p.text = "Characteristics:"
    char_p.font.size = Pt(22)
    char_p.font.bold = True
    char_p.font.color.rgb = BASE_COLOR
    
    characteristics = [
        "• Bitter taste",
        "• Soapy feel",
        "• Turn red litmus paper BLUE",
        "• Conduct electricity",
        "• React with acids"
    ]
    
    char_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(4), Inches(3.5))
    char_content_frame = char_content.text_frame
    char_content_frame.word_wrap = True
    
    for i, char in enumerate(characteristics):
        if i > 0:
            char_content_frame.add_paragraph()
        p = char_content_frame.paragraphs[i]
        p.text = char
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
    
    # Examples
    ex_title = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.3), Inches(0.4))
    ex_frame = ex_title.text_frame
    ex_p = ex_frame.paragraphs[0]
    ex_p.text = "Common Examples:"
    ex_p.font.size = Pt(22)
    ex_p.font.bold = True
    ex_p.font.color.rgb = BASE_COLOR
    
    examples = [
        "• NaOH - Drain cleaner",
        "• KOH - Soap making",
        "• Ca(OH)₂ - Slaked lime",
        "• NH₃ - Cleaning agent",
        "• Mg(OH)₂ - Milk of magnesia"
    ]
    
    ex_content = slide.shapes.add_textbox(Inches(5.2), Inches(2.5), Inches(4.3), Inches(3.5))
    ex_content_frame = ex_content.text_frame
    ex_content_frame.word_wrap = True
    
    for i, ex in enumerate(examples):
        if i > 0:
            ex_content_frame.add_paragraph()
        p = ex_content_frame.paragraphs[i]
        p.text = ex
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

def add_comparison_slide(prs):
    """Slide 5: Acids vs Bases Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 240, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Acids vs Bases - Quick Comparison"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = NEUTRAL_COLOR
    
    # Add table
    rows, cols = 8, 3
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.8)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table_shape.columns[0].width = Inches(2.5)
    table_shape.columns[1].width = Inches(3.25)
    table_shape.columns[2].width = Inches(3.25)
    
    # Header row
    headers = ["Property", "Acids", "Bases"]
    for i, header in enumerate(headers):
        cell = table_shape.cell(0, i)
        cell.text = header
        cell.fill.solid()
        if i == 0:
            cell.fill.fore_color.rgb = NEUTRAL_COLOR
        elif i == 1:
            cell.fill.fore_color.rgb = ACID_COLOR
        else:
            cell.fill.fore_color.rgb = BASE_COLOR
        
        text_frame = cell.text_frame
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.color.rgb = WHITE
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data rows
    data = [
        ["Taste", "Sour", "Bitter"],
        ["Feel", "Sharp", "Slippery/Soapy"],
        ["Litmus Paper", "Blue → Red", "Red → Blue"],
        ["Ions", "H⁺ ions", "OH⁻ ions"],
        ["Conductivity", "Good", "Good"],
        ["pH Value", "< 7", "> 7"],
        ["Reaction with metals", "Produces H₂", "Few metals"]
    ]
    
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
            text_frame = cell.text_frame
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            if col_idx == 0:
                text_frame.paragraphs[0].font.bold = True

def add_salts_slide(prs):
    """Slide 6: Salts - Definition & Formation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 230, 250)  # Light purple
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Salts - Definition & Formation"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = SALT_COLOR
    
    # Definition
    def_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.7))
    def_frame = def_box.text_frame
    def_frame.word_wrap = True
    def_p = def_frame.paragraphs[0]
    def_p.text = "Definition: Ionic compounds formed by the reaction between an acid and a base"
    def_p.font.size = Pt(18)
    def_p.font.bold = True
    def_p.font.color.rgb = SALT_COLOR
    
    # Neutralization Reaction
    rxn_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(0.3))
    rxn_frame = rxn_title.text_frame
    rxn_p = rxn_frame.paragraphs[0]
    rxn_p.text = "General Neutralization Reaction:"
    rxn_p.font.size = Pt(20)
    rxn_p.font.bold = True
    rxn_p.font.color.rgb = DARK_GRAY
    
    # Reaction equation
    eq_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.5))
    eq_frame = eq_box.text_frame
    eq_p = eq_frame.paragraphs[0]
    eq_p.text = "Acid + Base → Salt + Water"
    eq_p.font.size = Pt(26)
    eq_p.font.bold = True
    eq_p.font.color.rgb = RGBColor(200, 50, 50)
    eq_p.alignment = PP_ALIGN.CENTER
    
    # Examples
    ex_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.3))
    ex_frame = ex_title.text_frame
    ex_p = ex_frame.paragraphs[0]
    ex_p.text = "Examples of Salt Formation:"
    ex_p.font.size = Pt(20)
    ex_p.font.bold = True
    ex_p.font.color.rgb = DARK_GRAY
    
    examples = [
        "1. HCl + NaOH → NaCl + H₂O",
        "2. H₂SO₄ + 2NaOH → Na₂SO₄ + 2H₂O",
        "3. HNO₃ + NH₄OH → NH₄NO₃ + H₂O"
    ]
    
    ex_content = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.6), Inches(1.5))
    ex_content_frame = ex_content.text_frame
    ex_content_frame.word_wrap = True
    
    for i, ex in enumerate(examples):
        if i > 0:
            ex_content_frame.add_paragraph()
        p = ex_content_frame.paragraphs[i]
        p.text = ex
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Common salts
    salt_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.3))
    salt_frame = salt_title.text_frame
    salt_p = salt_frame.paragraphs[0]
    salt_p.text = "Common Salts in Daily Life:"
    salt_p.font.size = Pt(20)
    salt_p.font.bold = True
    salt_p.font.color.rgb = DARK_GRAY
    
    salts = [
        "🧂 NaCl (Table salt)  |  KNO₃ (Fertilizer)  |  Na₂CO₃ (Washing soda)",
        "🧪 CaSO₄ (Plaster)  |  KCl (Salt substitute)  |  NH₄NO₃ (Fertilizer)"
    ]
    
    salt_content = slide.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(8.6), Inches(1.3))
    salt_content_frame = salt_content.text_frame
    salt_content_frame.word_wrap = True
    
    for i, salt in enumerate(salts):
        if i > 0:
            salt_content_frame.add_paragraph()
        p = salt_content_frame.paragraphs[i]
        p.text = salt
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

def add_ph_scale_slide(prs):
    """Slide 7: The pH Scale"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "The pH Scale"
    title_p.font.size = Pt(50)
    title_p.font.bold = True
    title_p.font.color.rgb = NEUTRAL_COLOR
    
    # Definition
    def_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.7))
    def_frame = def_box.text_frame
    def_frame.word_wrap = True
    def_p = def_frame.paragraphs[0]
    def_p.text = "pH measures the concentration of H⁺ ions | Range: 0 - 14"
    def_p.font.size = Pt(18)
    def_p.font.bold = True
    def_p.font.color.rgb = DARK_GRAY
    
    # pH Scale Bar
    scale_y = Inches(1.95)
    scale_height = Inches(0.4)
    
    # Create color gradient boxes for pH scale
    pH_ranges = [
        (0, 1, "0 Highly\nAcidic", RGBColor(255, 50, 50)),
        (1, 3, "1-3", RGBColor(255, 100, 50)),
        (3, 5, "3-5", RGBColor(255, 150, 50)),
        (5, 7, "5-7\nAcidic", RGBColor(255, 200, 100)),
        (7, 7.5, "7\nNeutral", RGBColor(200, 200, 200)),
        (7.5, 9, "7-9\nBasic", RGBColor(100, 200, 255)),
        (9, 11, "9-11", RGBColor(50, 150, 255)),
        (11, 14, "11-14 Highly\nBasic", RGBColor(50, 100, 255)),
    ]
    
    for start, end, label, color in pH_ranges:
        x_pos = Inches(0.5 + (start / 14) * 9)
        width = Inches((end - start) / 14 * 9)
        
        rect = slide.shapes.add_shape(1, x_pos, scale_y, width, scale_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.color.rgb = DARK_GRAY
    
    # pH interpretation
    interp_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.6))
    interp_frame = interp_box.text_frame
    interp_frame.word_wrap = True
    interp_p = interp_frame.paragraphs[0]
    interp_p.text = "pH < 7: Acidic  |  pH = 7: Neutral  |  pH > 7: Basic/Alkaline"
    interp_p.font.size = Pt(18)
    interp_p.font.bold = True
    interp_p.font.color.rgb = DARK_GRAY
    interp_p.alignment = PP_ALIGN.CENTER
    
    # Common substances
    sub_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(0.3))
    sub_frame = sub_title.text_frame
    sub_p = sub_frame.paragraphs[0]
    sub_p.text = "pH Values of Common Substances:"
    sub_p.font.size = Pt(20)
    sub_p.font.bold = True
    sub_p.font.color.rgb = DARK_GRAY
    
    substances = [
        "Battery acid (0) | Gastric juice (1-2) | Lemon juice (2-3) | Vinegar (3-4) | Coffee (5)",
        "Milk (6-6.5) | Pure water (7) | Seawater (8) | Baking soda (8-9) | Ammonia (11-12)",
        "Drain cleaner (14)"
    ]
    
    sub_content = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(3))
    sub_content_frame = sub_content.text_frame
    sub_content_frame.word_wrap = True
    
    for i, sub in enumerate(substances):
        if i > 0:
            sub_content_frame.add_paragraph()
        p = sub_content_frame.paragraphs[i]
        p.text = sub
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

def add_indicators_slide(prs):
    """Slide 8: pH Indicators & Detection"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 250, 240)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "pH Indicators & Detection"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = NEUTRAL_COLOR
    
    # Types section
    types_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.3))
    types_frame = types_title.text_frame
    types_p = types_frame.paragraphs[0]
    types_p.text = "Types of Indicators:"
    types_p.font.size = Pt(20)
    types_p.font.bold = True
    types_p.font.color.rgb = DARK_GRAY
    
    # Indicator types table
    rows, cols = 5, 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(3)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table_shape.columns[0].width = Inches(2.5)
    table_shape.columns[1].width = Inches(3.25)
    table_shape.columns[2].width = Inches(3.25)
    
    # Headers
    headers = ["Indicator", "Acid Color", "Base Color"]
    for i, header in enumerate(headers):
        cell = table_shape.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEUTRAL_COLOR
        
        text_frame = cell.text_frame
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.color.rgb = WHITE
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data
    indicator_data = [
        ["Litmus Paper", "Red", "Blue"],
        ["Methyl Orange", "Red", "Yellow"],
        ["Phenolphthalein", "Colorless", "Pink"],
        ["Red Cabbage", "Pink/Red", "Green/Yellow"]
    ]
    
    for row_idx, row_data in enumerate(indicator_data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text
            
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
            text_frame = cell.text_frame
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            if col_idx == 0:
                text_frame.paragraphs[0].font.bold = True
    
    # Note
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    note_frame = note_box.text_frame
    note_frame.word_wrap = True
    
    note_text = [
        "✓ Litmus: Most common, quick indicator",
        "✓ Universal Indicator: Shows complete pH range (0-14)",
        "✓ pH Meter: Digital device for accurate measurement"
    ]
    
    for i, note in enumerate(note_text):
        if i > 0:
            note_frame.add_paragraph()
        p = note_frame.paragraphs[i]
        p.text = note
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

def add_neutralization_slide(prs):
    """Slide 9: Neutralization Reaction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 240, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Neutralization Reaction"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = SALT_COLOR
    
    # Definition
    def_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.7))
    def_frame = def_box.text_frame
    def_frame.word_wrap = True
    def_p = def_frame.paragraphs[0]
    def_p.text = "Definition: Reaction between acid and base producing salt and water. Ionic equation: H⁺(aq) + OH⁻(aq) → H₂O(l)"
    def_p.font.size = Pt(16)
    def_p.font.bold = True
    def_p.font.color.rgb = DARK_GRAY
    
    # Types of salts
    types_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(0.3))
    types_frame = types_title.text_frame
    types_p = types_frame.paragraphs[0]
    types_p.text = "Types of Salts Formed:"
    types_p.font.size = Pt(20)
    types_p.font.bold = True
    types_p.font.color.rgb = DARK_GRAY
    
    # Three columns for salt types
    salt_types = [
        ("Normal Salt", "Complete neutralization\nExample: HCl + NaOH\npH = 7", RGBColor(200, 200, 200)),
        ("Acidic Salt", "Excess acid used\nExample: HCl + NH₃\npH < 7", RGBColor(255, 150, 150)),
        ("Basic Salt", "Excess base used\nExample: NaOH + H₂CO₃\npH > 7", RGBColor(150, 200, 255))
    ]
    
    for idx, (title, desc, color) in enumerate(salt_types):
        x_pos = Inches(0.6 + idx * 3.1)
        
        # Box
        box_shape = slide.shapes.add_shape(1, x_pos, Inches(2.35), Inches(2.9), Inches(2.8))
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = color
        box_shape.line.color.rgb = DARK_GRAY
        box_shape.line.width = Pt(2)
        
        # Title
        title_box = slide.shapes.add_textbox(x_pos + Inches(0.1), Inches(2.5), Inches(2.7), Inches(0.4))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(16)
        title_p.font.bold = True
        title_p.font.color.rgb = DARK_GRAY
        title_p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos + Inches(0.15), Inches(3.05), Inches(2.6), Inches(2))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(13)
        desc_p.font.color.rgb = DARK_GRAY
        desc_p.alignment = PP_ALIGN.CENTER
    
    # Practical examples
    prac_title = slide.shapes.add_textbox(Inches(0.5), Inches(5.35), Inches(9), Inches(0.3))
    prac_frame = prac_title.text_frame
    prac_p = prac_frame.paragraphs[0]
    prac_p.text = "Real-Life Applications:"
    prac_p.font.size = Pt(18)
    prac_p.font.bold = True
    prac_p.font.color.rgb = DARK_GRAY
    
    apps = [
        "🏥 Treating heartburn: Stomach acid neutralized by antacids",
        "🌾 Agriculture: Acidic soil treated with lime (CaO) to raise pH",
        "🔧 Industrial: Wastewater treatment by neutralizing before discharge"
    ]
    
    app_content = slide.shapes.add_textbox(Inches(0.7), Inches(5.75), Inches(8.6), Inches(1.4))
    app_content_frame = app_content.text_frame
    app_content_frame.word_wrap = True
    
    for i, app in enumerate(apps):
        if i > 0:
            app_content_frame.add_paragraph()
        p = app_content_frame.paragraphs[i]
        p.text = app
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

def add_strong_weak_acids_slide(prs):
    """Slide 10: Strong & Weak Acids"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 235, 235)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Strong & Weak Acids"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = ACID_COLOR
    
    # Two columns
    strong_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.5), Inches(0.4))
    strong_frame = strong_title.text_frame
    strong_p = strong_frame.paragraphs[0]
    strong_p.text = "STRONG ACIDS"
    strong_p.font.size = Pt(22)
    strong_p.font.bold = True
    strong_p.font.color.rgb = RGBColor(220, 50, 50)
    
    weak_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.15), Inches(4.3), Inches(0.4))
    weak_frame = weak_title.text_frame
    weak_p = weak_frame.paragraphs[0]
    weak_p.text = "WEAK ACIDS"
    weak_p.font.size = Pt(22)
    weak_p.font.bold = True
    weak_p.font.color.rgb = RGBColor(200, 100, 100)
    
    # Strong acids content
    strong_acids = [
        "• Complete ionization (100%)",
        "• HCl - Hydrochloric acid",
        "• H₂SO₄ - Sulfuric acid",
        "• HNO₃ - Nitric acid",
        "• HBr - Hydrobromic acid",
        "• HClO₄ - Perchloric acid"
    ]
    
    strong_content = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(4.3), Inches(3.2))
    strong_content_frame = strong_content.text_frame
    strong_content_frame.word_wrap = True
    
    for i, acid in enumerate(strong_acids):
        if i > 0:
            strong_content_frame.add_paragraph()
        p = strong_content_frame.paragraphs[i]
        p.text = acid
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)
    
    # Weak acids content
    weak_acids = [
        "• Partial ionization (<100%)",
        "• CH₃COOH - Acetic acid",
        "• H₂CO₃ - Carbonic acid",
        "• H₃PO₄ - Phosphoric acid",
        "• C₆H₈O₇ - Citric acid",
        "• HCOOH - Formic acid"
    ]
    
    weak_content = slide.shapes.add_textbox(Inches(5.3), Inches(1.7), Inches(4.2), Inches(3.2))
    weak_content_frame = weak_content.text_frame
    weak_content_frame.word_wrap = True
    
    for i, acid in enumerate(weak_acids):
        if i > 0:
            weak_content_frame.add_paragraph()
        p = weak_content_frame.paragraphs[i]
        p.text = acid
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)
    
    # Comparison at bottom
    comp_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(2))
    comp_frame = comp_box.text_frame
    comp_frame.word_wrap = True
    
    comparisons = [
        "Strong acids conduct electricity excellently and have very low pH (1-2)",
        "Weak acids conduct electricity poorly and have higher pH (3-6)",
        "Example: Strong HCl = pH 1, Weak acetic acid = pH 3"
    ]
    
    for i, comp in enumerate(comparisons):
        if i > 0:
            comp_frame.add_paragraph()
        p = comp_frame.paragraphs[i]
        p.text = comp
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

def add_strong_weak_bases_slide(prs):
    """Slide 11: Strong & Weak Bases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 245, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Strong & Weak Bases"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = BASE_COLOR
    
    # Two columns
    strong_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(4.5), Inches(0.4))
    strong_frame = strong_title.text_frame
    strong_p = strong_frame.paragraphs[0]
    strong_p.text = "STRONG BASES"
    strong_p.font.size = Pt(22)
    strong_p.font.bold = True
    strong_p.font.color.rgb = RGBColor(50, 100, 220)
    
    weak_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.15), Inches(4.3), Inches(0.4))
    weak_frame = weak_title.text_frame
    weak_p = weak_frame.paragraphs[0]
    weak_p.text = "WEAK BASES"
    weak_p.font.size = Pt(22)
    weak_p.font.bold = True
    weak_p.font.color.rgb = RGBColor(100, 150, 220)
    
    # Strong bases content
    strong_bases = [
        "• Complete ionization (100%)",
        "• NaOH - Sodium hydroxide",
        "• KOH - Potassium hydroxide",
        "• Ca(OH)₂ - Calcium hydroxide",
        "• Ba(OH)₂ - Barium hydroxide",
        "• LiOH - Lithium hydroxide"
    ]
    
    strong_content = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(4.3), Inches(3.2))
    strong_content_frame = strong_content.text_frame
    strong_content_frame.word_wrap = True
    
    for i, base in enumerate(strong_bases):
        if i > 0:
            strong_content_frame.add_paragraph()
        p = strong_content_frame.paragraphs[i]
        p.text = base
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)
    
    # Weak bases content
    weak_bases = [
        "• Partial ionization (<100%)",
        "• NH₃ - Ammonia",
        "• CH₃NH₂ - Methylamine",
        "• Na₂CO₃ - Sodium carbonate",
        "(Produces OH⁻ through hydrolysis)"
    ]
    
    weak_content = slide.shapes.add_textbox(Inches(5.3), Inches(1.7), Inches(4.2), Inches(3.2))
    weak_content_frame = weak_content.text_frame
    weak_content_frame.word_wrap = True
    
    for i, base in enumerate(weak_bases):
        if i > 0:
            weak_content_frame.add_paragraph()
        p = weak_content_frame.paragraphs[i]
        p.text = base
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)
    
    # Comparison at bottom
    comp_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(2))
    comp_frame = comp_box.text_frame
    comp_frame.word_wrap = True
    
    comparisons = [
        "Strong bases conduct electricity excellently and have very high pH (12-14)",
        "Weak bases conduct electricity poorly and have lower pH (8-11)",
        "Example: Strong NaOH = pH 14, Weak ammonia = pH 11"
    ]
    
    for i, comp in enumerate(comparisons):
        if i > 0:
            comp_frame.add_paragraph()
        p = comp_frame.paragraphs[i]
        p.text = comp
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

def add_acid_reactions_slide(prs):
    """Slide 12: Chemical Properties of Acids"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 240, 240)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Chemical Properties of Acids"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = ACID_COLOR
    
    reactions = [
        ("1. Reaction with Metals", "Acid + Metal → Salt + H₂ gas\n2HCl + Zn → ZnCl₂ + H₂↑"),
        ("2. Reaction with Metal Oxides", "Acid + Oxide → Salt + Water\nH₂SO₄ + CuO → CuSO₄ + H₂O"),
        ("3. Reaction with Carbonates", "Acid + Carbonate → Salt + H₂O + CO₂↑\n2HCl + CaCO₃ → CaCl₂ + H₂O + CO₂↑"),
        ("4. Reaction with Bicarbonates", "Acid + Bicarbonate → Salt + H₂O + CO₂↑\nHCl + NaHCO₃ → NaCl + H₂O + CO₂↑"),
        ("5. Reaction with Bases", "Acid + Base → Salt + Water\nHCl + NaOH → NaCl + H₂O"),
        ("6. Reaction with Ammonia", "Acid + NH₃ → Ammonium salt\nHCl + NH₃ → NH₄Cl")
    ]
    
    y_start = Inches(1.1)
    for idx, (title, reaction) in enumerate(reactions):
        y_pos = y_start + Inches(idx * 0.95)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), y_pos, Inches(4), Inches(0.3))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(13)
        title_p.font.bold = True
        title_p.font.color.rgb = ACID_COLOR
        
        # Reaction
        rxn_box = slide.shapes.add_textbox(Inches(5), y_pos, Inches(4.5), Inches(0.6))
        rxn_frame = rxn_box.text_frame
        rxn_frame.word_wrap = True
        rxn_p = rxn_frame.paragraphs[0]
        rxn_p.text = reaction
        rxn_p.font.size = Pt(12)
        rxn_p.font.color.rgb = DARK_GRAY

def add_base_reactions_slide(prs):
    """Slide 13: Chemical Properties of Bases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 250, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Chemical Properties of Bases"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = BASE_COLOR
    
    reactions = [
        ("1. Reaction with Acids", "Base + Acid → Salt + Water\n2NaOH + H₂SO₄ → Na₂SO₄ + 2H₂O"),
        ("2. Reaction with Non-metal Oxides", "Base + CO₂ → Salt + Water\n2NaOH + CO₂ → Na₂CO₃ + H₂O"),
        ("3. Reaction with Certain Metals", "Base + Al → Salt + H₂ gas\n2NaOH + 2Al + 2H₂O → 2NaAlO₂ + 3H₂↑"),
        ("4. Reaction with Ammonium Salts", "Base + NH₄⁺ → NH₃↑ + Salt + Water\nNaOH + NH₄Cl → NaCl + NH₃↑ + H₂O"),
        ("5. Reaction with Non-metal Chlorides", "Base + Chloride → Oxides + Salt + Water\n3NaOH + PCl₃ → Na₃PO₃ + 3NaCl + H₂O")
    ]
    
    y_start = Inches(1.1)
    for idx, (title, reaction) in enumerate(reactions):
        y_pos = y_start + Inches(idx * 1.15)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), y_pos, Inches(4), Inches(0.3))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(13)
        title_p.font.bold = True
        title_p.font.color.rgb = BASE_COLOR
        
        # Reaction
        rxn_box = slide.shapes.add_textbox(Inches(5), y_pos, Inches(4.5), Inches(0.8))
        rxn_frame = rxn_box.text_frame
        rxn_frame.word_wrap = True
        rxn_p = rxn_frame.paragraphs[0]
        rxn_p.text = reaction
        rxn_p.font.size = Pt(12)
        rxn_p.font.color.rgb = DARK_GRAY

def add_salt_classification_slide(prs):
    """Slide 14: Salts - Classification & Properties"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 240, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Salt Classification"
    title_p.font.size = Pt(50)
    title_p.font.bold = True
    title_p.font.color.rgb = SALT_COLOR
    
    # Classification types
    types = [
        ("Normal Salts", "Complete neutralization\nNo H⁺ or OH⁻ remaining\npH = 7\nExample: NaCl", RGBColor(200, 200, 200)),
        ("Acidic Salts", "Weak base + Strong acid\nContains H⁺ ions\npH < 7\nExample: NH₄Cl", RGBColor(255, 150, 150)),
        ("Basic Salts", "Strong base + Weak acid\nContains OH⁻ groups\npH > 7\nExample: NaHCO₃", RGBColor(150, 200, 255))
    ]
    
    for idx, (title, desc, color) in enumerate(types):
        x_pos = Inches(0.6 + idx * 3)
        y_pos = Inches(1.2)
        
        # Box
        box_shape = slide.shapes.add_shape(1, x_pos, y_pos, Inches(2.8), Inches(3.2))
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = color
        box_shape.line.color.rgb = DARK_GRAY
        box_shape.line.width = Pt(2)
        
        # Title
        title_box = slide.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(0.15), Inches(2.6), Inches(0.4))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(16)
        title_p.font.bold = True
        title_p.font.color.rgb = DARK_GRAY
        title_p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.7), Inches(2.5), Inches(2.3))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(13)
        desc_p.font.color.rgb = DARK_GRAY
        desc_p.alignment = PP_ALIGN.CENTER
    
    # Other types
    other_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9), Inches(0.3))
    other_frame = other_title.text_frame
    other_p = other_frame.paragraphs[0]
    other_p.text = "Other Salt Types:"
    other_p.font.size = Pt(18)
    other_p.font.bold = True
    other_p.font.color.rgb = DARK_GRAY
    
    other_types = [
        "🔹 Double Salts: Contain two cations or anions (e.g., Alum - K₂SO₄·Al₂(SO₄)₃·24H₂O)",
        "🔹 Complex Salts: Contain complex ions (e.g., K₃[Fe(CN)₆])"
    ]
    
    other_content = slide.shapes.add_textbox(Inches(0.7), Inches(5.15), Inches(8.6), Inches(1.8))
    other_content_frame = other_content.text_frame
    other_content_frame.word_wrap = True
    
    for i, other in enumerate(other_types):
        if i > 0:
            other_content_frame.add_paragraph()
        p = other_content_frame.paragraphs[i]
        p.text = other
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

def add_common_salts_slide(prs):
    """Slide 15: Common Salts in Daily Life"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 245, 240)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Common Salts in Daily Life"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = NEUTRAL_COLOR
    
    # Create table
    rows, cols = 6, 4
    left = Inches(0.3)
    top = Inches(1.1)
    width = Inches(9.4)
    height = Inches(5.8)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table_shape.columns[0].width = Inches(1.8)
    table_shape.columns[1].width = Inches(2.2)
    table_shape.columns[2].width = Inches(2.7)
    table_shape.columns[3].width = Inches(2.7)
    
    # Headers
    headers = ["Salt Name", "Formula", "Uses", "Properties"]
    for i, header in enumerate(headers):
        cell = table_shape.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEUTRAL_COLOR
        
        text_frame = cell.text_frame
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(13)
        text_frame.paragraphs[0].font.color.rgb = WHITE
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data
    salt_data = [
        ["Table Salt", "NaCl", "Seasoning, preservation", "White crystals, soluble"],
        ["Baking Soda", "NaHCO₃", "Baking, cleaning, antacid", "White powder, basic"],
        ["Washing Soda", "Na₂CO₃", "Cleaning, glass making", "White powder, basic"],
        ["Fertilizer", "KNO₃ / NH₄NO₃", "Crop nutrition", "Crystals, soluble"],
        ["Epsom Salt", "MgSO₄", "Laxative, muscle relaxant", "White crystals, soluble"]
    ]
    
    for row_idx, row_data in enumerate(salt_data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text
            
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
            text_frame = cell.text_frame
            text_frame.paragraphs[0].font.size = Pt(11)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.word_wrap = True

def add_safety_slide(prs):
    """Slide 16: Hazards & Safety Precautions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 240, 240)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "⚠️  Hazards & Safety Precautions"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(200, 50, 50)
    
    # Two columns
    acid_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.5), Inches(0.3))
    acid_frame = acid_title.text_frame
    acid_p = acid_frame.paragraphs[0]
    acid_p.text = "⚠️ ACIDS - Safety Rules"
    acid_p.font.size = Pt(18)
    acid_p.font.bold = True
    acid_p.font.color.rgb = ACID_COLOR
    
    acids_rules = [
        "✗ Never add water to acid!",
        "✓ ALWAYS add Acid to Water (A to W)",
        "✓ Wear: Goggles, gloves, lab coat",
        "✓ Use in well-ventilated areas",
        "✓ Never touch concentrated acids",
        "✓ Use glass rods to stir"
    ]
    
    acid_content = slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(4.3), Inches(3.2))
    acid_content_frame = acid_content.text_frame
    acid_content_frame.word_wrap = True
    
    for i, rule in enumerate(acids_rules):
        if i > 0:
            acid_content_frame.add_paragraph()
        p = acid_content_frame.paragraphs[i]
        p.text = rule
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)
    
    base_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.3))
    base_frame = base_title.text_frame
    base_p = base_frame.paragraphs[0]
    base_p.text = "⚠️ BASES - Safety Rules"
    base_p.font.size = Pt(18)
    base_p.font.bold = True
    base_p.font.color.rgb = BASE_COLOR
    
    bases_rules = [
        "✓ Wear: Goggles, gloves, lab coat",
        "✓ Use in ventilated areas",
        "✓ Prevent splashing on skin",
        "✗ Never taste or ingest",
        "✓ Store in proper containers",
        "✓ Handle with care (slippery)"
    ]
    
    base_content = slide.shapes.add_textbox(Inches(5.3), Inches(1.55), Inches(4.2), Inches(3.2))
    base_content_frame = base_content.text_frame
    base_content_frame.word_wrap = True
    
    for i, rule in enumerate(bases_rules):
        if i > 0:
            base_content_frame.add_paragraph()
        p = base_content_frame.paragraphs[i]
        p.text = rule
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)
    
    # Note at bottom
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(1.8))
    note_frame = note_box.text_frame
    note_frame.word_wrap = True
    note_p = note_frame.paragraphs[0]
    note_p.text = "🔴 EMERGENCY: In case of acid/base spill on skin → Wash immediately with plenty of water for at least 10 minutes. Seek medical help if severe."
    note_p.font.size = Pt(14)
    note_p.font.bold = True
    note_p.font.color.rgb = RGBColor(200, 50, 50)
    note_p.alignment = PP_ALIGN.CENTER

def add_summary_slide(prs):
    """Slide 17: Summary & Key Points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Summary - Key Points to Remember"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = NEUTRAL_COLOR
    
    # Three columns
    key_points = [
        ("ACIDS", [
            "• Produce H⁺ ions",
            "• Turn blue litmus RED",
            "• Sour taste",
            "• pH < 7",
            "• React with bases"
        ], ACID_COLOR),
        ("BASES", [
            "• Produce OH⁻ ions",
            "• Turn red litmus BLUE",
            "• Bitter taste, soapy",
            "• pH > 7",
            "• React with acids"
        ], BASE_COLOR),
        ("SALTS", [
            "• From acid + base",
            "• Can be acidic/basic",
            "• Crystalline solids",
            "• In daily use",
            "• pH varies"
        ], SALT_COLOR)
    ]
    
    for idx, (title, points, color) in enumerate(key_points):
        x_pos = Inches(0.5 + idx * 3.1)
        
        # Title
        title_box = slide.shapes.add_textbox(x_pos, Inches(1.15), Inches(3), Inches(0.35))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(20)
        title_p.font.bold = True
        title_p.font.color.rgb = color
        title_p.alignment = PP_ALIGN.CENTER
        
        # Points
        points_box = slide.shapes.add_textbox(x_pos + Inches(0.1), Inches(1.65), Inches(2.8), Inches(4))
        points_frame = points_box.text_frame
        points_frame.word_wrap = True
        
        for i, point in enumerate(points):
            if i > 0:
                points_frame.add_paragraph()
            p = points_frame.paragraphs[i]
            p.text = point
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)

def add_misconceptions_slide(prs):
    """Slide 18: Common Misconceptions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 245, 240)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Common Misconceptions Clarified"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = NEUTRAL_COLOR
    
    misconceptions = [
        ("WRONG:", "pH 0 means no H⁺ ions", "RIGHT:", "pH 0 means VERY HIGH H⁺ concentration"),
        ("WRONG:", "All bases are toxic", "RIGHT:", "Some bases like NaHCO₃ are safe to eat"),
        ("WRONG:", "Strong acids are always concentrated", "RIGHT:", "Strong = complete ionization, Concentrated = amount dissolved"),
        ("WRONG:", "Neutral solutions are not reactive", "RIGHT:", "Neutral (pH=7) can still be reactive (hydrolysis)")
    ]
    
    y_start = Inches(1.2)
    for idx, (wrong_label, wrong_text, right_label, right_text) in enumerate(misconceptions):
        y_pos = y_start + Inches(idx * 1.3)
        
        # Wrong (Red background)
        wrong_box = slide.shapes.add_shape(1, Inches(0.5), y_pos, Inches(4.5), Inches(0.9))
        wrong_box.fill.solid()
        wrong_box.fill.fore_color.rgb = RGBColor(255, 200, 200)
        wrong_box.line.color.rgb = RGBColor(200, 50, 50)
        wrong_box.line.width = Pt(2)
        
        # Wrong text
        wrong_text_box = slide.shapes.add_textbox(Inches(0.6), y_pos + Inches(0.05), Inches(4.3), Inches(0.8))
        wrong_frame = wrong_text_box.text_frame
        wrong_frame.word_wrap = True
        wrong_p = wrong_frame.paragraphs[0]
        wrong_p.text = f"✗ {wrong_label}\n{wrong_text}"
        wrong_p.font.size = Pt(11)
        wrong_p.font.color.rgb = RGBColor(150, 0, 0)
        wrong_p.font.bold = True
        
        # Right (Green background)
        right_box = slide.shapes.add_shape(1, Inches(5.2), y_pos, Inches(4.3), Inches(0.9))
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(200, 255, 200)
        right_box.line.color.rgb = RGBColor(50, 150, 50)
        right_box.line.width = Pt(2)
        
        # Right text
        right_text_box = slide.shapes.add_textbox(Inches(5.3), y_pos + Inches(0.05), Inches(4.1), Inches(0.8))
        right_frame = right_text_box.text_frame
        right_frame.word_wrap = True
        right_p = right_frame.paragraphs[0]
        right_p.text = f"✓ {right_label}\n{right_text}"
        right_p.font.size = Pt(11)
        right_p.font.color.rgb = RGBColor(0, 100, 0)
        right_p.font.bold = True

def add_practice_slide(prs):
    """Slide 19: Practice Questions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 250, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Practice Questions"
    title_p.font.size = Pt(50)
    title_p.font.bold = True
    title_p.font.color.rgb = NEUTRAL_COLOR
    
    questions = [
        "1. Which of the following is a weak acid?\n   a) HCl  b) H₂SO₄  c) CH₃COOH ✓  d) HNO₃",
        "2. The pH of pure water at 25°C is:\n   a) 0  b) 7 ✓  c) 14  d) Undefined",
        "3. Balance: HCl + Ca(OH)₂ → CaCl₂ + H₂O\n   Answer: 2HCl + Ca(OH)₂ → CaCl₂ + 2H₂O",
        "4. What is produced when acid reacts with a carbonate?\n   Answer: Salt, Water, and Carbon Dioxide gas",
        "5. Name a real-world use of neutralization:\n   Answer: Treating heartburn with antacids"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(8.8), Inches(5.8))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, question in enumerate(questions):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = question
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        p.level = 0

def add_closing_slide(prs):
    """Slide 20: Closing Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 80, 150)
    
    # Main message
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    msg_frame = msg_box.text_frame
    msg_frame.word_wrap = True
    msg_p = msg_frame.paragraphs[0]
    msg_p.text = "Acids and Bases are EVERYWHERE!"
    msg_p.font.size = Pt(60)
    msg_p.font.bold = True
    msg_p.font.color.rgb = WHITE
    msg_p.alignment = PP_ALIGN.CENTER
    
    # Subtext
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    sub_frame = sub_box.text_frame
    sub_frame.word_wrap = True
    sub_p = sub_frame.paragraphs[0]
    sub_p.text = "From food we eat to medicines we take,\nUnderstanding their properties helps us understand chemistry and life!"
    sub_p.font.size = Pt(24)
    sub_p.font.italic = True
    sub_p.font.color.rgb = RGBColor(200, 220, 255)
    sub_p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.7))
    footer_frame = footer_box.text_frame
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "NCERT Class 10 Chemistry Complete Course"
    footer_p.font.size = Pt(18)
    footer_p.font.color.rgb = RGBColor(180, 200, 255)
    footer_p.alignment = PP_ALIGN.CENTER

def main():
    """Generate the complete presentation"""
    print("🎬 Generating Acids, Bases & Salts Presentation...")
    print("=" * 60)
    
    prs = create_presentation()
    
    slides_data = [
        ("Slide 1: Title", add_title_slide),
        ("Slide 2: Learning Objectives", add_objectives_slide),
        ("Slide 3: What are Acids?", add_acids_slide),
        ("Slide 4: What are Bases?", add_bases_slide),
        ("Slide 5: Acids vs Bases Comparison", add_comparison_slide),
        ("Slide 6: Salts - Definition & Formation", add_salts_slide),
        ("Slide 7: The pH Scale", add_ph_scale_slide),
        ("Slide 8: pH Indicators & Detection", add_indicators_slide),
        ("Slide 9: Neutralization Reaction", add_neutralization_slide),
        ("Slide 10: Strong & Weak Acids", add_strong_weak_acids_slide),
        ("Slide 11: Strong & Weak Bases", add_strong_weak_bases_slide),
        ("Slide 12: Chemical Properties of Acids", add_acid_reactions_slide),
        ("Slide 13: Chemical Properties of Bases", add_base_reactions_slide),
        ("Slide 14: Salt Classification", add_salt_classification_slide),
        ("Slide 15: Common Salts in Daily Life", add_common_salts_slide),
        ("Slide 16: Hazards & Safety", add_safety_slide),
        ("Slide 17: Summary", add_summary_slide),
        ("Slide 18: Common Misconceptions", add_misconceptions_slide),
        ("Slide 19: Practice Questions", add_practice_slide),
        ("Slide 20: Closing", add_closing_slide),
    ]
    
    for slide_name, slide_func in slides_data:
        try:
            slide_func(prs)
            print(f"✓ {slide_name}")
        except Exception as e:
            print(f"✗ {slide_name}: {str(e)}")
    
    # Save presentation
    output_filename = "Acids_Bases_Salts_Class10.pptx"
    prs.save(output_filename)
    
    print("=" * 60)
    print(f"✅ Presentation created successfully!")
    print(f"📄 File: {output_filename}")
    print(f"📊 Total Slides: 20")
    print(f"🎨 Color-coded: Acids (Red), Bases (Blue), Salts (Purple)")
    print(f"📚 NCERT Class 10 Chemistry - Complete Coverage")
    print("\nPresentation is ready for download and viewing!")

if __name__ == "__main__":
    main()
